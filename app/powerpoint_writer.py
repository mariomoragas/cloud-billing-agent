from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

FONT_NAME = "Verdana"
TITLE_FONT = "Georgia"
APP_ROOT = Path(__file__).resolve().parent
PPT_TEMPLATE_DIR = APP_ROOT / "assets" / "ppt_template"
ASSET_ORACLE_LOGO = PPT_TEMPLATE_DIR / "oracle_logo.png"
ASSET_COVER_BG = PPT_TEMPLATE_DIR / "cover_bg.png"
ASSET_TOP_RIGHT = PPT_TEMPLATE_DIR / "top_right.png"
ASSET_BOTTOM_RIGHT = PPT_TEMPLATE_DIR / "bottom_right.png"
COLOR_BG = RGBColor(247, 244, 242)
COLOR_PANEL = RGBColor(255, 255, 255)
COLOR_BORDER = RGBColor(229, 214, 206)
COLOR_PRIMARY = RGBColor(199, 70, 52)
COLOR_PRIMARY_DARK = RGBColor(142, 47, 35)
COLOR_SECONDARY = RGBColor(232, 184, 175)
COLOR_TEXT = RGBColor(27, 27, 27)
COLOR_MUTED = RGBColor(95, 91, 87)
COLOR_DARK = RGBColor(27, 27, 27)
FOOTER_TEXT = "Copyright © 2026, Oracle | Confidential: Restricted"


def write_powerpoint_report(
    output_path: Path,
    raw_df: pd.DataFrame,
    service_summary_df: pd.DataFrame,
    region_summary_df: pd.DataFrame,
    oci_mapping_df: pd.DataFrame,
    llm_report_df: pd.DataFrame | None = None,
    llm_migration_df: pd.DataFrame | None = None,
    llm_recommendations_df: pd.DataFrame | None = None,
    llm_confidence_df: pd.DataFrame | None = None,
    report_name: str = "",
    company_name: str = "",
    project_name: str = "",
) -> None:
    output_path.parent.mkdir(parents=True, exist_ok=True)

    presentation = Presentation()
    _configure_presentation(presentation)
    slide_number = 1
    _add_title_slide(presentation, raw_df, report_name, company_name, project_name)
    slide_number += 1
    _add_section_slide(
        presentation,
        title="Executive Summary",
        subtitle="Financial highlights and OCI conversion overview",
        slide_number=slide_number,
    )
    slide_number += 1
    _add_kpi_slide(presentation, raw_df, service_summary_df, oci_mapping_df, slide_number=slide_number)
    slide_number += 1
    _add_top_services_slide(presentation, service_summary_df, slide_number=slide_number)
    slide_number += 1
    _add_service_share_slide(presentation, service_summary_df, slide_number=slide_number)
    slide_number += 1
    _add_region_slide(presentation, region_summary_df, slide_number=slide_number)
    slide_number += 1
    _add_migration_complexity_slide(presentation, oci_mapping_df, slide_number=slide_number)
    slide_number += 1
    _add_mapping_slide(presentation, oci_mapping_df, slide_number=slide_number)
    slide_number += 1
    if llm_report_df is not None and not llm_report_df.empty:
        _add_section_slide(
            presentation,
            title="LLM FinOps Analysis",
            subtitle="Baseline, OCI projection, ROI and migration plan synthesized by LLM",
            slide_number=slide_number,
        )
        slide_number += 1
        _add_llm_overview_slide(
            presentation,
            llm_report_df=llm_report_df,
            llm_confidence_df=llm_confidence_df,
            slide_number=slide_number,
        )
        slide_number += 1
        _add_llm_plan_slide(
            presentation,
            llm_migration_df=llm_migration_df,
            llm_recommendations_df=llm_recommendations_df,
            slide_number=slide_number,
        )
        slide_number += 1
    _add_section_slide(
        presentation,
        title="OCI Review Focus",
        subtitle="Services requiring manual validation before conversion",
        slide_number=slide_number,
    )
    slide_number += 1
    _add_unmapped_slide(presentation, oci_mapping_df, slide_number=slide_number)
    presentation.save(output_path)


def _configure_presentation(presentation: Presentation) -> None:
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)


def _add_title_slide(
    presentation: Presentation,
    raw_df: pd.DataFrame,
    report_name: str,
    company_name: str,
    project_name: str,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_slide(slide, dark=True)
    _add_picture_if_exists(slide, ASSET_COVER_BG, 0, 0, 13.333, 7.5)
    _add_picture_if_exists(slide, ASSET_ORACLE_LOGO, 0.76, 1.02, 1.67, 0.35)
    header_project = project_name.strip() or "Cloud Billing Conversion Executive Summary"
    _add_title(slide, header_project, 0.76, 2.45, 7.4, dark=True)
    client_name = company_name.strip() or _derive_client_name(report_name)
    if client_name:
        client_box = slide.shapes.add_textbox(Inches(0.76), Inches(1.65), Inches(6.0), Inches(0.45))
        _configure_text_frame(client_box.text_frame, margin_left=0.02, margin_right=0.02)
        client_paragraph = client_box.text_frame.paragraphs[0]
        client_paragraph.text = client_name
        _style_paragraph(client_paragraph, bold=False, size=13, color=(255, 255, 255))
    period = _mode_or_default(raw_df, "period", "N/A")
    cloud = _mode_or_default(raw_df, "cloud", "cloud").upper()
    subtitle = slide.shapes.add_textbox(Inches(0.76), Inches(3.85), Inches(7.6), Inches(0.8))
    _configure_text_frame(subtitle.text_frame, margin_left=0.02, margin_right=0.02)
    paragraph = subtitle.text_frame.paragraphs[0]
    paragraph.text = f"Assessment scope | Billing period: {period} | Source: {cloud}"
    _style_paragraph(paragraph, bold=False, size=20, color=_rgb_tuple(COLOR_SECONDARY))

    total_cost = float(raw_df["cost"].sum()) if not raw_df.empty else 0.0
    currency = _mode_or_default(raw_df, "currency", "USD")
    _add_hero_metric(slide, "Total analyzed cost", f"{total_cost:,.2f} {currency}", 0.76, 4.55, dark=True)
    _add_footer(slide, FOOTER_TEXT, dark=True, slide_number=1)


def _add_section_slide(
    presentation: Presentation,
    *,
    title: str,
    subtitle: str,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_slide(slide, dark=True)
    _add_picture_if_exists(slide, ASSET_COVER_BG, 0, 0, 13.333, 7.5)
    _add_picture_if_exists(slide, ASSET_ORACLE_LOGO, 0.76, 1.02, 1.67, 0.35)
    _add_title(slide, title, 0.76, 2.55, 6.6, dark=True)
    sub_box = slide.shapes.add_textbox(Inches(0.76), Inches(3.45), Inches(6.8), Inches(0.5))
    _configure_text_frame(sub_box.text_frame, margin_left=0.02, margin_right=0.02)
    sub_paragraph = sub_box.text_frame.paragraphs[0]
    sub_paragraph.text = subtitle
    _style_paragraph(sub_paragraph, bold=False, size=18, color=_rgb_tuple(COLOR_SECONDARY))
    _add_footer(slide, FOOTER_TEXT, dark=True, slide_number=slide_number)


def _add_kpi_slide(
    presentation: Presentation,
    raw_df: pd.DataFrame,
    service_summary_df: pd.DataFrame,
    oci_mapping_df: pd.DataFrame,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_slide(slide)
    _add_standard_frame(slide, "Key Metrics", slide_number)

    total_cost = float(raw_df["cost"].sum()) if not raw_df.empty else 0.0
    currency = _mode_or_default(raw_df, "currency", "USD")
    service_count = int(raw_df["service_name_original"].nunique()) if not raw_df.empty else 0
    unmapped_count = (
        int((oci_mapping_df["oci_service"] == "REVIEW_REQUIRED").sum())
        if not oci_mapping_df.empty
        else 0
    )
    cards = [
        ("Total Cost", f"{total_cost:,.2f} {currency}"),
        ("Distinct Services", str(service_count)),
        ("Unmapped Services", str(unmapped_count)),
    ]
    for index, (label, value) in enumerate(cards):
        _add_kpi_card(slide, label, value, 0.76 + (index * 3.0), 1.6)

    service_label_column = _resolve_service_label_column(service_summary_df)
    top_service = _top_label(service_summary_df, service_label_column)
    insight = (
        f"The current period is concentrated in {top_service}. "
        f"Use the next slides to prioritize OCI conversion, service rationalization and financial review."
    )
    _add_insight_box(slide, insight, 0.76, 4.25, 9.45, 1.25)


def _add_top_services_slide(
    presentation: Presentation,
    service_summary_df: pd.DataFrame,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_slide(slide)
    _add_standard_frame(slide, "Top Services by Cost", slide_number)

    service_label_column = _resolve_service_label_column(service_summary_df)
    top_services = _top_with_others(service_summary_df, service_label_column, "total_cost", 8)
    chart_data = ChartData()
    chart_data.categories = list(top_services["label"])
    chart_data.add_series("Cost", list(top_services["value"]))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.76),
        Inches(1.55),
        Inches(8.7),
        Inches(4.8),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.number_format = '#,##0'
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = COLOR_PRIMARY
    chart.series[0].format.line.color.rgb = COLOR_PRIMARY
    _set_chart_fonts(chart)
    _add_subtitle(slide, "Monthly spend concentration across the main billed services", 0.76, 1.18, 6.8)
    _add_side_note(
        slide,
        "Prioritize OCI mapping and optimization on the first bars. They usually capture most of the monthly spend.",
        9.65,
        1.7,
        2.8,
        2.4,
    )


def _add_service_share_slide(
    presentation: Presentation,
    service_summary_df: pd.DataFrame,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_slide(slide)
    _add_standard_frame(slide, "Service Cost Share", slide_number)

    service_label_column = _resolve_service_label_column(service_summary_df)
    top_services = _top_with_others(service_summary_df, service_label_column, "total_cost", 5)
    chart_data = ChartData()
    chart_data.categories = list(top_services["label"])
    chart_data.add_series("Cost Share", list(top_services["value"]))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(0.76),
        Inches(1.5),
        Inches(7.6),
        Inches(4.9),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.show_percentage = True
    chart.plots[0].data_labels.show_category_name = True
    _set_chart_fonts(chart)
    _add_subtitle(slide, "Share of the total cost represented by the leading services", 0.76, 1.18, 6.8)
    _add_side_note(
        slide,
        "Use this view to explain concentration. If one service dominates, map and optimize it first for OCI.",
        9.05,
        1.75,
        2.95,
        2.25,
    )


def _add_region_slide(
    presentation: Presentation,
    region_summary_df: pd.DataFrame,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_slide(slide)
    _add_standard_frame(slide, "Top Regions by Cost", slide_number)

    top_regions = _top_with_others(region_summary_df, "region", "total_cost", 6)
    top_regions["label"] = top_regions["label"].replace("", "UNSPECIFIED")
    chart_data = ChartData()
    chart_data.categories = list(top_regions["label"])
    chart_data.add_series("Cost", list(top_regions["value"]))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.76),
        Inches(1.55),
        Inches(8.7),
        Inches(4.8),
        chart_data,
    ).chart
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.number_format = '#,##0'
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = COLOR_SECONDARY
    chart.series[0].format.line.color.rgb = COLOR_SECONDARY
    _set_chart_fonts(chart)
    _add_subtitle(slide, "Regional distribution of the billed workload after the applied filters", 0.76, 1.18, 7.2)
    _add_side_note(
        slide,
        "Regional concentration highlights potential latency, network and data residency impacts before OCI redesign.",
        9.65,
        1.75,
        2.8,
        2.25,
    )


def _add_migration_complexity_slide(
    presentation: Presentation,
    oci_mapping_df: pd.DataFrame,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_slide(slide)
    _add_standard_frame(slide, "Complexidade de Migracao", slide_number)
    _add_subtitle(slide, "Top servicos por custo posicionados por score de complexidade", 0.76, 1.18, 7.8)

    points = _build_complexity_points(oci_mapping_df, top_n=10)
    if points.empty:
        _add_insight_box(
            slide,
            "Nao ha dados suficientes de complexidade para os servicos mapeados.",
            0.76,
            2.0,
            8.8,
            1.6,
            title_size=12,
            body_size=12,
        )
        return

    panel_left = 0.76
    panel_top = 1.55
    panel_width = 8.7
    panel_height = 4.9

    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(panel_left),
        Inches(panel_top),
        Inches(panel_width),
        Inches(panel_height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(240, 238, 236)
    panel.line.color.rgb = COLOR_BORDER

    vertical = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(panel_left + (panel_width * 0.5)),
        Inches(panel_top),
        Inches(0.01),
        Inches(panel_height),
    )
    vertical.fill.solid()
    vertical.fill.fore_color.rgb = COLOR_PRIMARY
    vertical.line.fill.background()

    horizontal = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(panel_left),
        Inches(panel_top + (panel_height * 0.5)),
        Inches(panel_width),
        Inches(0.01),
    )
    horizontal.fill.solid()
    horizontal.fill.fore_color.rgb = COLOR_PRIMARY
    horizontal.line.fill.background()

    x_axis = slide.shapes.add_textbox(Inches(panel_left + 2.4), Inches(panel_top + panel_height + 0.05), Inches(4.0), Inches(0.25))
    _configure_text_frame(x_axis.text_frame, margin_left=0.01, margin_right=0.01)
    p_x = x_axis.text_frame.paragraphs[0]
    p_x.text = "Prazo da Migracao"
    _style_paragraph(p_x, bold=False, size=12, color=_rgb_tuple(COLOR_MUTED))

    y_axis = slide.shapes.add_textbox(Inches(panel_left - 0.65), Inches(panel_top + 1.6), Inches(0.6), Inches(2.0))
    _configure_text_frame(y_axis.text_frame, margin_left=0.01, margin_right=0.01)
    p_y = y_axis.text_frame.paragraphs[0]
    p_y.text = "Custo da Migracao"
    _style_paragraph(p_y, bold=False, size=11, color=_rgb_tuple(COLOR_MUTED))

    for index, row in points.iterrows():
        rank = float(index + 1)
        complexity = float(row["complexity_score"])
        label = str(row["label"])
        score_color = _complexity_color(int(complexity))

        x_ratio = (rank - 1.0) / max(len(points) - 1, 1)
        y_ratio = (complexity - 1.0) / 4.0
        x = panel_left + 0.35 + (x_ratio * (panel_width - 1.1))
        y = panel_top + (panel_height - 0.75) - (y_ratio * (panel_height - 1.0))
        y = y - (0.18 if (index % 2) else 0.0)

        bubble = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL,
            Inches(x),
            Inches(y),
            Inches(0.42),
            Inches(0.42),
        )
        bubble.fill.solid()
        bubble.fill.fore_color.rgb = score_color
        bubble.line.fill.background()
        bubble_tf = bubble.text_frame
        _configure_text_frame(bubble_tf, margin_left=0.0, margin_right=0.0, margin_top=0.0, margin_bottom=0.0)
        b = bubble_tf.paragraphs[0]
        b.text = str(int(complexity))
        _style_paragraph(b, bold=True, size=11, color=(255, 255, 255))

        tag = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(max(panel_left + 0.05, x - 0.92)),
            Inches(y + 0.44),
            Inches(1.85),
            Inches(0.36),
        )
        tag.fill.solid()
        tag.fill.fore_color.rgb = _complexity_tag_color(int(complexity))
        tag.line.fill.background()
        tag_tf = tag.text_frame
        _configure_text_frame(tag_tf, margin_left=0.03, margin_right=0.03)
        t = tag_tf.paragraphs[0]
        t.text = label[:30]
        _style_paragraph(t, bold=False, size=10, color=_rgb_tuple(COLOR_TEXT))

    low_complexity = points[points["complexity_score"] <= 2]
    low_count = int(len(low_complexity))
    total_top_cost = float(points["total_cost"].sum()) if not points.empty else 0.0
    low_cost_share = (
        float(low_complexity["total_cost"].sum()) / total_top_cost * 100.0
        if total_top_cost > 0
        else 0.0
    )
    cloud = _mode_or_default(oci_mapping_df, "cloud", "cloud").upper()

    right = slide.shapes.add_textbox(Inches(9.85), Inches(2.1), Inches(3.1), Inches(3.2))
    _configure_text_frame(right.text_frame, margin_left=0.01, margin_right=0.01)
    p1 = right.text_frame.paragraphs[0]
    p1.text = f"{low_count} dos 10"
    _style_paragraph(p1, bold=False, size=30, color=_rgb_tuple(COLOR_TEXT))
    p2 = right.text_frame.add_paragraph()
    p2.text = f"principais servicos na {cloud} possuem baixa complexidade de migracao,"
    _style_paragraph(p2, bold=False, size=13, color=_rgb_tuple(COLOR_TEXT))
    p3 = right.text_frame.add_paragraph()
    p3.text = f"representando {low_cost_share:.0f}% do custo total"
    _style_paragraph(p3, bold=True, size=16, color=_rgb_tuple(COLOR_TEXT))


def _add_unmapped_slide(
    presentation: Presentation,
    oci_mapping_df: pd.DataFrame,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_slide(slide)
    _add_standard_frame(slide, "Services Requiring OCI Review", slide_number)
    _add_subtitle(slide, "Manual validation queue before final OCI equivalence and sizing", 0.76, 1.18, 7.4)

    unmapped_df = (
        oci_mapping_df[oci_mapping_df["oci_service"] == "REVIEW_REQUIRED"]
        .sort_values("total_cost", ascending=False)
        .head(8)
    )

    _add_side_note(
        slide,
        "These services require manual equivalence validation before final OCI sizing or migration assumptions.",
        9.6,
        1.6,
        2.85,
        2.0,
    )

    box = slide.shapes.add_textbox(Inches(0.82), Inches(1.75), Inches(8.35), Inches(3.25))
    frame = box.text_frame
    _configure_text_frame(frame)

    if unmapped_df.empty:
        paragraph = frame.paragraphs[0]
        paragraph.text = "All services in this report have an initial OCI mapping."
        _style_paragraph(paragraph, bold=False, size=14, color=_rgb_tuple(COLOR_TEXT))
        return

    first = True
    for _, row in unmapped_df.iterrows():
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        cost = float(row["total_cost"]) if "total_cost" in row else 0.0
        paragraph.text = f"{row['service_name_original']} | {cost:,.2f} {row.get('primary_currency', 'USD')}"
        _style_paragraph(paragraph, bold=False, size=14, color=_rgb_tuple(COLOR_TEXT))


def _add_mapping_slide(
    presentation: Presentation,
    oci_mapping_df: pd.DataFrame,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_slide(slide)
    _add_standard_frame(slide, "AWS to OCI Mapping (Consolidated)", slide_number)
    _add_subtitle(
        slide,
        "Consolidated pairs from Mapeamento_OCI (source service -> OCI service)",
        0.76,
        1.18,
        8.6,
    )

    mapping_df = _consolidate_mapping_pairs(oci_mapping_df)
    _add_side_note(
        slide,
        "Repeated lines are grouped by identical source and destination.",
        9.6,
        1.6,
        2.85,
        2.0,
    )

    box = slide.shapes.add_textbox(Inches(0.82), Inches(1.75), Inches(8.35), Inches(5.15))
    frame = box.text_frame
    _configure_text_frame(frame)

    if mapping_df.empty:
        paragraph = frame.paragraphs[0]
        paragraph.text = "No mapped services found for this report."
        _style_paragraph(paragraph, bold=False, size=18, color=_rgb_tuple(COLOR_TEXT))
        return

    max_lines = 12
    visible = mapping_df.head(max_lines)
    first = True
    for _, row in visible.iterrows():
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        paragraph.text = f"{row['service_name_original']} --> {row['oci_service']}"
        _style_paragraph(paragraph, bold=False, size=11, color=_rgb_tuple(COLOR_TEXT))

    remaining = len(mapping_df) - len(visible)
    if remaining > 0:
        paragraph = frame.add_paragraph()
        paragraph.text = f"+ {remaining} additional consolidated mappings not shown."
        _style_paragraph(paragraph, bold=True, size=11, color=_rgb_tuple(COLOR_MUTED))


def _add_llm_overview_slide(
    presentation: Presentation,
    llm_report_df: pd.DataFrame,
    llm_confidence_df: pd.DataFrame | None,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_slide(slide)
    _add_standard_frame(slide, "LLM Baseline, Projection and ROI", slide_number)
    _add_subtitle(slide, "Numbers and assumptions generated from billing + OCI mapping context", 0.76, 1.18, 8.0)

    current_total = _llm_value(llm_report_df, "baseline", "current_total_cost")
    base_total = _llm_value(llm_report_df, "projection", "base_total")
    savings_pct = _llm_value(llm_report_df, "savings", "base_savings_pct")
    roi_pct = _llm_value(llm_report_df, "business_case", "roi_pct")
    payback = _llm_value(llm_report_df, "business_case", "payback_months")
    currency = str(_llm_value(llm_report_df, "baseline", "currency", default="USD"))
    analysis_mode = str(_llm_value(llm_report_df, "meta", "analysis_mode", default="unknown"))

    cards = [
        ("Baseline", f"{current_total:,.2f} {currency}" if isinstance(current_total, (int, float)) else str(current_total)),
        ("OCI Base", f"{base_total:,.2f} {currency}" if isinstance(base_total, (int, float)) else str(base_total)),
        ("Savings Base", f"{savings_pct:.2f}%" if isinstance(savings_pct, (int, float)) else str(savings_pct)),
    ]
    for index, (label, value) in enumerate(cards):
        _add_kpi_card(slide, label, value, 0.76 + (index * 2.95), 1.6)

    summary = str(_llm_value(llm_report_df, "summary", "executive_summary", default="")).strip()
    if not summary:
        summary = "No executive summary returned."
    _add_insight_box(slide, summary, 0.76, 4.0, 8.7, 1.7, title_size=11, body_size=11)

    confidence_note = "No confidence information."
    if llm_confidence_df is not None and not llm_confidence_df.empty:
        preview = []
        for _, row in llm_confidence_df.head(3).iterrows():
            preview.append(f"{row.get('topic', '')}: {row.get('level', '')}")
        confidence_note = "\n".join(preview)
    _add_side_note(
        slide,
        f"ROI: {_format_percentage(roi_pct)}\nPayback (months): {_format_months(payback)} Mode: {analysis_mode}\n{confidence_note}",
        9.65,
        1.75,
        2.75,
        3.9,
        body_size=12,
    )


def _add_llm_plan_slide(
    presentation: Presentation,
    llm_migration_df: pd.DataFrame | None,
    llm_recommendations_df: pd.DataFrame | None,
    slide_number: int,
) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    _paint_slide(slide)
    _add_standard_frame(slide, "LLM Migration Plan and Recommendations", slide_number)
    _add_subtitle(slide, "Phases, risks, dependencies and OCI architecture recommendations", 0.76, 1.18, 8.4)

    left_box = slide.shapes.add_textbox(Inches(0.82), Inches(1.75), Inches(11.31), Inches(1.45))
    left_frame = left_box.text_frame
    _configure_text_frame(left_frame)

    lines: list[str] = []
    if llm_migration_df is not None and not llm_migration_df.empty:
        for _, row in llm_migration_df.head(4).iterrows():
            phase = str(row.get("phase", "")).strip()
            duration = str(row.get("duration", "")).strip()
            quick = str(row.get("quick_wins", "")).strip()
            lines.append(f"{phase} ({duration}) - {quick}")
    if not lines:
        lines.append("No migration plan returned.")

    first = True
    for line in lines:
        paragraph = left_frame.paragraphs[0] if first else left_frame.add_paragraph()
        first = False
        paragraph.text = line
        _style_paragraph(paragraph, bold=False, size=10, color=_rgb_tuple(COLOR_TEXT))

    right_note = []
    if llm_recommendations_df is not None and not llm_recommendations_df.empty:
        for _, row in llm_recommendations_df.head(7).iterrows():
            right_note.append(f"- {str(row.get('item', '')).strip()}")
    if not right_note:
        right_note = ["- No recommendations returned."]
    _add_side_note(
        slide,
        "Recommendations\n" + "\n".join(right_note),
        0.76,
        3.31,
        11.64,
        3.64,
        title_size=9,
        body_size=9,
    )


def _paint_slide(slide, dark: bool = False) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_DARK if dark else COLOR_BG
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)
    if not dark:
        _add_picture_if_exists(slide, ASSET_TOP_RIGHT, 8.6, 0.0, 4.73, 1.28)
        _add_picture_if_exists(slide, ASSET_BOTTOM_RIGHT, 10.7, 6.1, 2.2, 0.9)


def _add_oracle_banner(slide, dark: bool = False) -> None:
    banner = slide.shapes.add_textbox(Inches(0.76), Inches(0.38), Inches(3.0), Inches(0.45))
    paragraph = banner.text_frame.paragraphs[0]
    paragraph.text = "ORACLE"
    _style_paragraph(
        paragraph,
        bold=True,
        size=12,
        color=(255, 255, 255) if dark else _rgb_tuple(COLOR_PRIMARY),
    )


def _add_standard_frame(slide, title: str, slide_number: int) -> None:
    _add_oracle_banner(slide)
    _add_title(slide, title, 0.76, 0.52, 8.7)
    _add_title_accent(slide, 0.76, 1.08)
    _add_footer(slide, FOOTER_TEXT, slide_number=slide_number)


def _add_title(slide, text: str, left: float, top: float, width: float, dark: bool = False) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.8))
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    color = (255, 255, 255) if dark else _rgb_tuple(COLOR_TEXT)
    paragraph.font.name = TITLE_FONT
    paragraph.font.size = Pt(24 if not dark else 28)
    paragraph.font.bold = False
    paragraph.font.color.rgb = RGBColor(*color)


def _add_title_accent(slide, left: float, top: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(0.42),
        Inches(0.05),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PRIMARY
    shape.line.fill.background()


def _add_hero_metric(slide, label: str, value: str, left: float, top: float, dark: bool = False) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(4.2),
        Inches(1.8),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255) if not dark else RGBColor(251, 249, 248)
    shape.line.color.rgb = COLOR_SECONDARY
    frame = shape.text_frame
    _configure_text_frame(frame)
    frame.clear()
    p1 = frame.paragraphs[0]
    p1.text = label
    _style_paragraph(p1, bold=False, size=16, color=_rgb_tuple(COLOR_MUTED))
    p2 = frame.add_paragraph()
    p2.text = value
    _style_paragraph(p2, bold=True, size=24, color=_rgb_tuple(COLOR_TEXT))


def _add_kpi_card(slide, label: str, value: str, left: float, top: float) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(2.6),
        Inches(2.0),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_PANEL
    shape.line.color.rgb = COLOR_BORDER

    frame = shape.text_frame
    _configure_text_frame(frame)
    frame.clear()
    p1 = frame.paragraphs[0]
    p1.text = label
    _style_paragraph(p1, bold=False, size=16, color=_rgb_tuple(COLOR_MUTED))
    p2 = frame.add_paragraph()
    p2.text = value
    _style_paragraph(p2, bold=True, size=24, color=_rgb_tuple(COLOR_TEXT))


def _add_subtitle(slide, text: str, left: float, top: float, width: float) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.35))
    _configure_text_frame(box.text_frame, margin_left=0.01, margin_right=0.01)
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    _style_paragraph(paragraph, bold=False, size=11, color=_rgb_tuple(COLOR_MUTED))


def _add_insight_box(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    title_size: int = 13,
    body_size: int = 16,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 248, 245)
    shape.line.color.rgb = COLOR_SECONDARY
    frame = shape.text_frame
    _configure_text_frame(frame)
    p1 = frame.paragraphs[0]
    p1.text = "Executive takeaway"
    _style_paragraph(p1, bold=True, size=title_size, color=_rgb_tuple(COLOR_PRIMARY_DARK))
    p2 = frame.add_paragraph()
    p2.text = text
    _style_paragraph(p2, bold=False, size=body_size, color=_rgb_tuple(COLOR_TEXT))


def _add_side_note(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    title_size: int = 13,
    body_size: int = 14,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = COLOR_BORDER
    frame = shape.text_frame
    _configure_text_frame(frame)
    p1 = frame.paragraphs[0]
    p1.text = "Oracle note"
    _style_paragraph(p1, bold=True, size=title_size, color=_rgb_tuple(COLOR_PRIMARY))
    for line in text.splitlines():
        p2 = frame.add_paragraph()
        p2.text = line
        _style_paragraph(p2, bold=False, size=body_size, color=_rgb_tuple(COLOR_TEXT))


def _add_footer(slide, text: str, dark: bool = False, slide_number: int | None = None) -> None:
    if not dark:
        line = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.76),
            Inches(7.03),
            Inches(5.1),
            Inches(0.015),
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(139, 133, 128)
        line.line.fill.background()
    shape = slide.shapes.add_textbox(Inches(1.12), Inches(6.95), Inches(5.2), Inches(0.22))
    _configure_text_frame(shape.text_frame, margin_left=0.01, margin_right=0.01)
    paragraph = shape.text_frame.paragraphs[0]
    paragraph.text = text
    color = _rgb_tuple(COLOR_SECONDARY) if dark else (139, 133, 128)
    _style_paragraph(paragraph, bold=False, size=8, color=color)
    if slide_number is not None:
        number_box = slide.shapes.add_textbox(Inches(0.76), Inches(6.92), Inches(0.3), Inches(0.25))
        _configure_text_frame(number_box.text_frame, margin_left=0.0, margin_right=0.0)
        number_paragraph = number_box.text_frame.paragraphs[0]
        number_paragraph.text = str(slide_number)
        _style_paragraph(number_paragraph, bold=False, size=8, color=color)


def _style_paragraph(paragraph, *, bold: bool, size: int, color: tuple[int, int, int]) -> None:
    paragraph.font.name = FONT_NAME
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(*color)


def _set_chart_fonts(chart) -> None:
    if chart.has_title and chart.chart_title is not None:
        for paragraph in chart.chart_title.text_frame.paragraphs:
            paragraph.font.name = FONT_NAME
            paragraph.font.bold = True
            paragraph.font.size = Pt(12)
    if chart.has_legend and chart.legend is not None:
        chart.legend.font.name = FONT_NAME
        chart.legend.font.size = Pt(12)
    try:
        category_axis = chart.category_axis
        category_axis.tick_labels.font.name = FONT_NAME
        category_axis.tick_labels.font.size = Pt(12)
    except (AttributeError, ValueError):
        pass
    try:
        value_axis = chart.value_axis
        value_axis.tick_labels.font.name = FONT_NAME
        value_axis.tick_labels.font.size = Pt(12)
    except (AttributeError, ValueError):
        pass
    if chart.plots:
        try:
            labels = chart.plots[0].data_labels
            labels.font.name = FONT_NAME
            labels.font.size = Pt(12)
        except Exception:
            pass


def _top_with_others(
    df: pd.DataFrame,
    label_col: str,
    value_col: str,
    top_n: int,
) -> pd.DataFrame:
    if df.empty:
        return pd.DataFrame(columns=["label", "value"])
    working = df.copy()
    working[label_col] = working[label_col].fillna("").astype(str).replace("", "UNSPECIFIED")
    top = working[[label_col, value_col]].head(top_n).copy()
    top.columns = ["label", "value"]
    others_total = float(working.iloc[top_n:][value_col].sum())
    if others_total > 0:
        top.loc[len(top)] = ["Others", others_total]
    return top


def _consolidate_mapping_pairs(oci_mapping_df: pd.DataFrame) -> pd.DataFrame:
    if oci_mapping_df.empty:
        return pd.DataFrame()

    working = oci_mapping_df.copy()
    working["service_name_original"] = (
        working["service_name_original"].fillna("").astype(str).str.strip()
    )
    working["oci_service"] = working["oci_service"].fillna("").astype(str).str.strip()
    working["primary_currency"] = (
        working.get("primary_currency", pd.Series(["USD"] * len(working)))
        .fillna("USD")
        .astype(str)
        .str.strip()
    )

    mapped = working[
        (working["service_name_original"] != "")
        & (working["oci_service"] != "")
        & (working["oci_service"] != "REVIEW_REQUIRED")
    ]
    if mapped.empty:
        return pd.DataFrame()

    grouped = (
        mapped.groupby(["service_name_original", "oci_service"], as_index=False)
        .agg(
            total_cost=("total_cost", "sum"),
            total_usage_quantity=("total_usage_quantity", "sum"),
            row_count=("oci_service", "count"),
            primary_currency=("primary_currency", _mode_text_or_default),
        )
        .sort_values("total_cost", ascending=False)
    )
    return grouped


def _mode_text_or_default(values: pd.Series) -> str:
    cleaned = values.fillna("").astype(str).str.strip()
    cleaned = cleaned[cleaned != ""]
    if cleaned.empty:
        return "USD"
    return str(cleaned.mode().iloc[0])


def _format_percentage(value: object) -> str:
    if isinstance(value, (int, float)):
        return f"{float(value):.2f}%"
    return "-"


def _format_months(value: object) -> str:
    if isinstance(value, (int, float)):
        return f"{float(value):.1f}"
    return "-"


def _llm_value(
    llm_report_df: pd.DataFrame,
    section: str,
    metric: str,
    default: object = 0.0,
) -> object:
    scoped = llm_report_df[
        (llm_report_df["section"] == section) & (llm_report_df["metric"] == metric)
    ]
    if scoped.empty:
        return default
    value = scoped.iloc[0]["value"]
    if isinstance(value, (int, float)):
        return value
    if value is None:
        return default
    text = str(value).strip()
    if text == "":
        return default
    try:
        return float(text)
    except ValueError:
        return text


def _top_label(df: pd.DataFrame, column: str) -> str:
    if df.empty or column not in df.columns:
        return "the leading service"
    value = str(df.iloc[0][column]).strip()
    return value or "the leading service"


def _resolve_service_label_column(service_summary_df: pd.DataFrame) -> str:
    if "chart_group_label" in service_summary_df.columns:
        values = service_summary_df["chart_group_label"].fillna("").astype(str).str.strip()
        if (values != "").any():
            return "chart_group_label"
    return "service_name_original"


def _build_complexity_points(oci_mapping_df: pd.DataFrame, top_n: int = 10) -> pd.DataFrame:
    if oci_mapping_df.empty:
        return pd.DataFrame(columns=["label", "total_cost", "complexity_score"])

    working = oci_mapping_df.copy()
    if "oci_service" in working.columns:
        working = working[working["oci_service"].astype(str).str.strip() != "REVIEW_REQUIRED"]
    if working.empty:
        return pd.DataFrame(columns=["label", "total_cost", "complexity_score"])

    service = working.get("service_name_original", pd.Series([""] * len(working))).fillna("").astype(str).str.strip()
    product = working.get("primary_product_code", pd.Series([""] * len(working))).fillna("").astype(str).str.strip()
    working["label"] = product.where(product != "", service).replace("", "UNSPECIFIED")
    working["total_cost"] = pd.to_numeric(working.get("total_cost", 0.0), errors="coerce").fillna(0.0)
    working["complexity_score"] = (
        pd.to_numeric(working.get("complexity_score", 2), errors="coerce")
        .fillna(2)
        .astype(int)
        .clip(1, 5)
    )

    grouped = (
        working.groupby(["label", "complexity_score"], as_index=False)
        .agg(total_cost=("total_cost", "sum"))
        .sort_values("total_cost", ascending=False)
        .head(top_n)
        .reset_index(drop=True)
    )
    return grouped


def _complexity_color(score: int) -> RGBColor:
    if score <= 2:
        return RGBColor(124, 151, 88)
    if score == 3:
        return RGBColor(187, 146, 72)
    return RGBColor(153, 73, 64)


def _complexity_tag_color(score: int) -> RGBColor:
    if score <= 2:
        return RGBColor(196, 219, 192)
    if score == 3:
        return RGBColor(242, 221, 161)
    return RGBColor(236, 202, 197)


def _rgb_tuple(color: RGBColor) -> tuple[int, int, int]:
    return (color[0], color[1], color[2])


def _add_picture_if_exists(slide, path: Path, left: float, top: float, width: float, height: float) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))


def _derive_client_name(report_name: str) -> str:
    cleaned = str(report_name or "").replace("_", " ").replace("-", " ").strip()
    cleaned = " ".join(cleaned.split())
    if not cleaned:
        return ""
    suffixes = ("report", "billing", "invoice", "aws", "azure", "gcp")
    parts = [part for part in cleaned.split() if part.lower() not in suffixes]
    candidate = " ".join(parts).strip()
    return candidate.title() if candidate else cleaned.title()


def _mode_or_default(df: pd.DataFrame, column: str, default: str) -> str:
    if df.empty or column not in df.columns:
        return default
    series = df[column].fillna("").astype(str).str.strip()
    series = series[series != ""]
    if series.empty:
        return default
    return str(series.mode().iloc[0])


def _configure_text_frame(
    frame,
    *,
    margin_left: float = 0.04,
    margin_right: float = 0.04,
    margin_top: float = 0.03,
    margin_bottom: float = 0.03,
) -> None:
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = Inches(margin_left)
    frame.margin_right = Inches(margin_right)
    frame.margin_top = Inches(margin_top)
    frame.margin_bottom = Inches(margin_bottom)
